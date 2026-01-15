import base64
import io
import os

from dash import Dash, dcc, html, Input, Output, State
import pandas as pd
from pptx import Presentation

OUTPUT_FILE = "output.xlsx"

# -------------------------------
# Backend logic
# -------------------------------

def extract_assessments_from_ppt(contents):
    """
    contents: base64-encoded file contents from Dash Upload component
    """
    content_type, content_string = contents.split(',')  # HTML file uploads come in 2 parts : data:<mime-type>; base64,<data>. We just want the data.
    decoded = base64.b64decode(content_string)

    prs = Presentation(io.BytesIO(decoded)) # from Presentation import
                                            # this reads ppt/presentation.xml, slides1.xml, ppt/media/* etc.. and 
                                            # Builds Python objects for Presentation, SLides, Shapes, Text frames, Layouts, Masters, Relationships
    assessments = []

    for slide_number, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue

            text = shape.text.strip()
            if text.startswith("Assessment"):
                assessments.append({
                    "Slide": slide_number,
                    "Assessment Text": text
                })

    return assessments


def write_to_excel(data):
    df = pd.DataFrame(data)
    df.to_excel(OUTPUT_FILE, index=False, engine="openpyxl")


# -------------------------------
# Dash App
# -------------------------------

app = Dash(__name__)

app.layout = html.Div(
    style={
        "width": "60%",
        "margin": "auto",
        "textAlign": "center",
        "fontFamily": "Arial"
    },
    children=[
        html.H2("PowerPoint → Excel Assessment Extractor"),

        dcc.Upload(
            id="upload-ppt",
            children=html.Div([
                "Drag and Drop or ",
                html.A("Select a PowerPoint File")  # !! When a file is dropped, the JS frontend(browser) encodes the file as 
                                                    #...base64 and stores it in the components 'content' property.  "contents" lives..
                                                    # .. in browser memory. If you want it to persist, you have to save it (excel, dataset,..)
            ]),
            style={
                "width": "100%",
                "height": "120px",
                "lineHeight": "120px",
                "borderWidth": "2px",
                "borderStyle": "dashed",
                "borderRadius": "10px",
                "textAlign": "center",
                "margin": "20px 0"
            },
            multiple=False,
            accept=".pptx"
        ),

        html.Div(id="status-message", style={"marginTop": "20px"}),

        html.Hr(),

        html.Div(id="preview-table")
    ]
)


# -------------------------------
# Callbacks
# -------------------------------

@app.callback(
    Output("status-message", "children"),  # See Notes at bottom for details on these callbacks.
    Output("preview-table", "children"),
    Input("upload-ppt", "contents"),        
    State("upload-ppt", "filename"),
    prevent_initial_call=True
)
def process_upload(contents, filename): # func is called from callback trigger.  Contents passes to 'extract_assess..."
    if contents is None:
        return "", ""

    if not filename.lower().endswith(".pptx"):
        return html.Div("Please upload a .pptx file.", style={"color": "red"}), ""

    try:
        data = extract_assessments_from_ppt(contents)  # func is called from callback trigger.  Contents passes to 'extract_assess..."

        if not data:
            return html.Div(
                "No assessments found in this presentation.",
                style={"color": "orange"}
            ), ""

        write_to_excel(data)  # *** FUTURE CHANGE *** Calls the function to write to excel. Prob want to write to Dataset in Palantir.

        df = pd.DataFrame(data)  # Need this here for the return below which prints the output in the browser.

        return (
            html.Div(
                f"Successfully extracted {len(data)} assessments. Excel file overwritten: {OUTPUT_FILE}",
                style={"color": "green"}
            ),
            dcc.Graph(  # simple table preview
                figure={
                    "data": [{
                        "type": "table",
                        "header": {
                            "values": list(df.columns),
                            "align": "left"
                        },
                        "cells": {
                            "values": [df[col] for col in df.columns],
                            "align": "left"
                        }
                    }],
                    "layout": {"margin": {"t": 10}}
                }
            )
        )

    except Exception as e:
        return html.Div(f"Error: {str(e)}", style={"color": "red"}), ""


# -------------------------------
# Run App
# -------------------------------

if __name__ == "__main__":
    app.run(debug=True)


#------------------------------------------
#       Callback information for clarity.
#------------------------------------------
# Input("upload-ppt", "contents")
# -- Dash comes with dcc.Upload button which we named "upload-ppt".  We'll use this dcc.Upload to drop our pptx files.

#    Output("status-message", "children"),  -Component of app.Layout = html.Div(stlyes[], children[html.Div, dcc.Upload(*ox to upload),...]
#                                           - This "Output" is for the text feedback to the user (1 upload 2 outputs)
#    Output("preview-table", "children"),   - This "Output" is for the preview table in the website (1 upload 2 outputs)
#    Input("upload-ppt", "contents"),        
#    State("upload-ppt", "filename"),
#    prevent_initial_call=True              -- This prevents the function from executing upon initial load 
